"""
core/services/parsers.py
~~~~~~~~~~~~~~~~~~~~~~~~~
Text extraction parsers for every supported file format.

Spec: PDF / DOCX / DOC / XLSX / XLS / CSV / TXT / HTML / RTF / PPTX
      + image files handled by TesseractEngine directly.

Design decisions:
- Each parser is a small function returning ParseResult (structured output).
- DocumentParser.parse() is the single entry point — dispatches by extension.
- For PDF: try text extraction first (pypdf); fall back to OCR flag if
  extracted text is below a meaningful threshold (< 50 chars per page).
- Metadata extraction is best-effort; parsers never raise on metadata failure.
"""

import io
import logging
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

MIN_PDF_TEXT_CHARS_PER_PAGE = 50  # below this → mark requires_ocr=True


@dataclass
class PageResult:
    page_number: int
    text: str
    word_count: int = 0

    def __post_init__(self):
        if not self.word_count:
            self.word_count = len(self.text.split())


@dataclass
class ParseResult:
    pages: list[PageResult] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    requires_ocr: bool = False
    error: str = ""
    success: bool = True

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages)

    @property
    def total_word_count(self) -> int:
        return sum(p.word_count for p in self.pages)

    @property
    def page_count(self) -> int:
        return len(self.pages)


# ---------------------------------------------------------------------------
# Individual parsers
# ---------------------------------------------------------------------------

def parse_pdf(file_obj) -> ParseResult:
    """Extract text from a native (text-layer) PDF using pypdf."""
    try:
        from pypdf import PdfReader
        file_obj.seek(0)
        reader = PdfReader(file_obj)
        pages = []
        total_chars = 0

        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            total_chars += len(text)
            pages.append(PageResult(page_number=i, text=text))

        # Decide if OCR is needed
        avg_chars = total_chars / len(reader.pages) if reader.pages else 0
        requires_ocr = avg_chars < MIN_PDF_TEXT_CHARS_PER_PAGE

        # Metadata
        meta = {}
        if reader.metadata:
            for k in ("/Title", "/Author", "/CreationDate", "/ModDate"):
                v = reader.metadata.get(k)
                if v:
                    meta[k.lstrip("/")] = str(v)

        return ParseResult(
            pages=pages,
            metadata=meta,
            requires_ocr=requires_ocr,
        )
    except Exception as e:
        logger.error("PDF parse error: %s", e)
        return ParseResult(success=False, error=str(e), requires_ocr=True)


def parse_docx(file_obj) -> ParseResult:
    """Extract text and metadata from DOCX files."""
    try:
        from docx import Document
        file_obj.seek(0)
        doc = Document(file_obj)
        # Treat the whole document as a single "page" with paragraph separators
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)

        # Also extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            full_text += "\n\n" + "\n".join(rows)

        meta = {}
        try:
            cp = doc.core_properties
            if cp.author:    meta["Author"] = cp.author
            if cp.title:     meta["Title"]  = cp.title
            if cp.created:   meta["CreationDate"] = str(cp.created)
            if cp.modified:  meta["ModDate"] = str(cp.modified)
        except Exception:
            pass

        return ParseResult(
            pages=[PageResult(page_number=1, text=full_text)],
            metadata=meta,
        )
    except Exception as e:
        logger.error("DOCX parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_doc(file_path: Path) -> ParseResult:
    """Extract text from legacy .DOC via docx2txt (requires antiword on Linux;
    on Windows we fall back to docx2txt which handles most modern .doc files)."""
    try:
        import docx2txt
        text = docx2txt.process(str(file_path))
        return ParseResult(pages=[PageResult(page_number=1, text=text or "")])
    except Exception as e:
        logger.error("DOC parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_xlsx(file_obj) -> ParseResult:
    """Convert XLSX to text: each sheet becomes a page."""
    try:
        import openpyxl
        file_obj.seek(0)
        wb = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)
        pages = []
        for sheet_idx, sheet in enumerate(wb.worksheets, start=1):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) if cell is not None else "" for cell in row
                )
                if row_text.strip():
                    rows.append(row_text)
            text = f"Sheet: {sheet.title}\n" + "\n".join(rows)
            pages.append(PageResult(page_number=sheet_idx, text=text))
        return ParseResult(pages=pages, metadata={"sheets": len(wb.worksheets)})
    except Exception as e:
        logger.error("XLSX parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_xls(file_path: Path) -> ParseResult:
    """Convert legacy XLS via xlrd."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(file_path))
        pages = []
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            rows = []
            for row_idx in range(sheet.nrows):
                row_text = " | ".join(str(sheet.cell_value(row_idx, col))
                                      for col in range(sheet.ncols))
                if row_text.strip():
                    rows.append(row_text)
            text = f"Sheet: {sheet.name}\n" + "\n".join(rows)
            pages.append(PageResult(page_number=sheet_idx + 1, text=text))
        return ParseResult(pages=pages)
    except Exception as e:
        logger.error("XLS parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_csv(file_obj) -> ParseResult:
    """Read CSV as structured text."""
    try:
        import csv
        file_obj.seek(0)
        # Detect encoding
        raw = file_obj.read()
        file_obj.seek(0)
        try:
            text_content = raw.decode("utf-8")
        except UnicodeDecodeError:
            text_content = raw.decode("latin-1")
        reader = csv.reader(io.StringIO(text_content))
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
        return ParseResult(pages=[PageResult(page_number=1, text="\n".join(rows))])
    except Exception as e:
        logger.error("CSV parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_txt(file_obj) -> ParseResult:
    try:
        file_obj.seek(0)
        raw = file_obj.read()
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("latin-1")
        return ParseResult(pages=[PageResult(page_number=1, text=text)])
    except Exception as e:
        logger.error("TXT parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_html(file_obj) -> ParseResult:
    try:
        from bs4 import BeautifulSoup
        file_obj.seek(0)
        raw = file_obj.read()
        try:
            html = raw.decode("utf-8")
        except UnicodeDecodeError:
            html = raw.decode("latin-1")
        soup = BeautifulSoup(html, "lxml")
        # Remove scripts and styles
        for tag in soup(["script", "style", "head", "nav", "footer"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        meta = {}
        title_tag = soup.find("title")
        if title_tag:
            meta["Title"] = title_tag.get_text()
        return ParseResult(
            pages=[PageResult(page_number=1, text=text)],
            metadata=meta,
        )
    except Exception as e:
        logger.error("HTML parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_rtf(file_obj) -> ParseResult:
    try:
        from striprtf.striprtf import rtf_to_text
        file_obj.seek(0)
        rtf_content = file_obj.read().decode("latin-1", errors="replace")
        text = rtf_to_text(rtf_content)
        return ParseResult(pages=[PageResult(page_number=1, text=text)])
    except Exception as e:
        logger.error("RTF parse error: %s", e)
        return ParseResult(success=False, error=str(e))


def parse_pptx(file_obj) -> ParseResult:
    """Extract text slide-by-slide from PPTX. Each slide = one page."""
    try:
        from pptx import Presentation
        file_obj.seek(0)
        prs = Presentation(file_obj)
        pages = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs)
                        if line.strip():
                            parts.append(line)
            text = "\n".join(parts)
            pages.append(PageResult(page_number=slide_idx, text=text))
        return ParseResult(pages=pages)
    except Exception as e:
        logger.error("PPTX parse error: %s", e)
        return ParseResult(success=False, error=str(e))


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

class DocumentParser:
    """
    Single entry point.  parse(file_obj, filename, file_path) dispatches
    to the correct parser by file extension.
    """

    def parse(
        self,
        file_obj,
        filename: str,
        file_path: Optional[Path] = None,
    ) -> ParseResult:
        ext = Path(filename).suffix.lower()
        logger.info("Parsing %s (ext=%s)", filename, ext)

        try:
            if ext == ".pdf":
                return parse_pdf(file_obj)
            elif ext == ".docx":
                return parse_docx(file_obj)
            elif ext == ".doc":
                return parse_doc(file_path or Path(filename))
            elif ext == ".xlsx":
                return parse_xlsx(file_obj)
            elif ext == ".xls":
                return parse_xls(file_path or Path(filename))
            elif ext == ".csv":
                return parse_csv(file_obj)
            elif ext in (".txt",):
                return parse_txt(file_obj)
            elif ext in (".html", ".htm"):
                return parse_html(file_obj)
            elif ext == ".rtf":
                return parse_rtf(file_obj)
            elif ext == ".pptx":
                return parse_pptx(file_obj)
            elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif"):
                # Image files: signal requires_ocr, no direct text extraction
                return ParseResult(requires_ocr=True, pages=[])
            else:
                return ParseResult(
                    success=False,
                    error=f"Unsupported extension: {ext}",
                )
        except Exception as e:
            logger.exception("Unexpected parse error for %s", filename)
            return ParseResult(success=False, error=str(e))
